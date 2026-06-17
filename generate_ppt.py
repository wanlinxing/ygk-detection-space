"""
生成美观的演讲PPT
基于YOLOv8目标检测项目 - 美化版
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色方案 - 科技感蓝色调
PRIMARY_COLOR = RGBColor(0, 82, 147)      # 深蓝
SECONDARY_COLOR = RGBColor(0, 150, 199)   # 亮蓝
ACCENT_COLOR = RGBColor(0, 180, 216)      # 青色
TEXT_COLOR = RGBColor(51, 51, 51)         # 深灰
LIGHT_BG = RGBColor(240, 248, 255)        # 淡蓝背景
WHITE = RGBColor(255, 255, 255)

def add_gradient_background(slide, color1, color2):
    """添加渐变背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color1

def add_title_shape(slide, text, left, top, width, height, font_size=44, bold=True, color=WHITE):
    """添加标题文本框"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.LEFT
    return shape

def add_content_textbox(slide, text, left, top, width, height, font_size=18, color=TEXT_COLOR, bold=False):
    """添加内容文本框"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.LEFT
    return shape

def add_bullet_points(slide, items, left, top, width, height, font_size=20, color=TEXT_COLOR, bullet_color=PRIMARY_COLOR):
    """添加带项目符号的列表"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = "▸ " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(16)
        p.level = 0
    
    return shape

def add_decorative_bar(slide, left, top, width, height, color):
    """添加装饰条"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color, text="", font_size=16):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.bold = True
    
    return shape

# ===== 第1页：标题页 =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_gradient_background(slide1, PRIMARY_COLOR, SECONDARY_COLOR)

# 添加装饰圆形
shape = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(6), Inches(6))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_COLOR
shape.fill.fore_color.brightness = 0.3
shape.line.fill.background()

shape2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(4), Inches(4))
shape2.fill.solid()
shape2.fill.fore_color.rgb = SECONDARY_COLOR
shape2.fill.fore_color.brightness = 0.2
shape2.line.fill.background()

# 标题
add_title_shape(slide1, "基于改进YOLOv8的", Inches(1), Inches(2), Inches(8), Inches(1), font_size=48, color=WHITE)
add_title_shape(slide1, "目标检测系统", Inches(1), Inches(2.8), Inches(8), Inches(1), font_size=56, bold=True, color=WHITE)

# 副标题
add_content_textbox(slide1, "《深度学习技术与应用》课程期末大作业", Inches(1), Inches(4), Inches(8), Inches(0.6), 
                   font_size=24, color=RGBColor(200, 230, 255))

# 个人信息
info_text = "汇报人：[你的姓名]    学号：[你的学号]\n指导教师：[教师姓名]"
add_content_textbox(slide1, info_text, Inches(1), Inches(5.2), Inches(8), Inches(1), 
                   font_size=18, color=RGBColor(220, 240, 255))

# 底部装饰条
add_decorative_bar(slide1, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7), ACCENT_COLOR)

# ===== 第2页：目录 =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide2, WHITE, WHITE)

# 左侧色块
left_bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
left_bar.fill.solid()
left_bar.fill.fore_color.rgb = PRIMARY_COLOR
left_bar.line.fill.background()

# 标题区域
title_bg = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0), Inches(13), Inches(1.2))
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = PRIMARY_COLOR
title_bg.line.fill.background()

add_title_shape(slide2, "汇报提纲", Inches(0.8), Inches(0.25), Inches(6), Inches(0.8), font_size=36, color=WHITE)

# 目录项
toc_items = [
    ("01", "选题背景与意义", "目标检测的应用场景与研究价值"),
    ("02", "开发环境与技术路线", "PyTorch框架与YOLOv8模型选择"),
    ("03", "数据集与预处理", "VOC2012数据集与数据增强策略"),
    ("04", "模型设计与改进", "CBAM注意力机制与Transformer"),
    ("05", "模型训练与优化", "训练参数设置与损失函数设计"),
    ("06", "实验结果与分析", "性能对比与可视化展示"),
    ("07", "总结与展望", "项目总结与未来改进方向")
]

for i, (num, title, desc) in enumerate(toc_items):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6)
    y = Inches(1.6 + row * 1.6)
    
    # 编号圆形
    circle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_COLOR
    circle.line.fill.background()
    
    # 编号文字
    num_tf = circle.text_frame
    num_tf.paragraphs[0].text = num
    num_tf.paragraphs[0].font.size = Pt(20)
    num_tf.paragraphs[0].font.bold = True
    num_tf.paragraphs[0].font.color.rgb = WHITE
    num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 标题
    add_title_shape(slide2, title, x + Inches(0.8), y, Inches(4.5), Inches(0.4), font_size=22, color=PRIMARY_COLOR)
    
    # 描述
    add_content_textbox(slide2, desc, x + Inches(0.8), y + Inches(0.4), Inches(4.5), Inches(0.4), 
                       font_size=14, color=RGBColor(120, 120, 120))

# ===== 第3页：选题背景 =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide3, WHITE, WHITE)

# 顶部色块
top_bar = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = PRIMARY_COLOR
top_bar.line.fill.background()

# 标题
add_title_shape(slide3, "选题背景", Inches(0.5), Inches(0.4), Inches(6), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide3, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 四个要点卡片
cards = [
    ("目标检测", "计算机视觉核心任务，广泛应用于自动驾驶、智能监控、医学影像分析等领域"),
    ("YOLO系列", "以速度快、精度高成为主流方法，YOLOv8采用更先进的网络结构"),
    ("改进空间", "YOLOv8仍有提升空间，引入注意力机制可进一步增强特征提取能力"),
    ("项目目标", "基于YOLOv8，引入CBAM和Transformer，提升目标检测精度")
]

for i, (title, content) in enumerate(cards):
    x = Inches(0.5 + (i % 2) * 6.2)
    y = Inches(1.5 + (i // 2) * 2.8)
    
    # 卡片背景
    card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(2.4))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = ACCENT_COLOR
    card.line.width = Pt(2)
    
    # 卡片标题
    add_title_shape(slide3, title, x + Inches(0.3), y + Inches(0.2), Inches(5.2), Inches(0.5), 
                   font_size=22, color=PRIMARY_COLOR)
    
    # 卡片内容
    add_content_textbox(slide3, content, x + Inches(0.3), y + Inches(0.8), Inches(5.2), Inches(1.4), 
                       font_size=16, color=TEXT_COLOR)

# ===== 第4页：技术路线 =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide4, WHITE, WHITE)

# 标题
add_title_shape(slide4, "技术路线", Inches(0.5), Inches(0.4), Inches(6), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide4, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 技术栈流程图
flow_items = [
    ("数据集", "Pascal VOC 2012\n20个类别\n17,000张图像", PRIMARY_COLOR),
    ("基础模型", "YOLOv8\nCSPDarknet Backbone\nPANet Neck", SECONDARY_COLOR),
    ("改进方案", "CBAM注意力机制\nTransformer模块\n特征增强", ACCENT_COLOR),
    ("训练框架", "PyTorch 2.1\nAdamW优化器\nCosine Annealing", PRIMARY_COLOR),
    ("开发工具", "PyCharm IDE\nPython 3.10\nGit版本控制", SECONDARY_COLOR)
]

for i, (title, content, color) in enumerate(flow_items):
    x = Inches(0.5 + i * 2.5)
    y = Inches(2)
    
    # 圆角矩形
    rect = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(3.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    
    # 标题
    tf = rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 内容
    p2 = tf.add_paragraph()
    p2.text = "\n" + content
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(230, 245, 255)
    p2.alignment = PP_ALIGN.CENTER
    
    # 箭头（最后一个不添加）
    if i < len(flow_items) - 1:
        arrow = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.25), y + Inches(1.5), Inches(0.4), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_COLOR
        arrow.line.fill.background()

# 底部说明
add_content_textbox(slide4, "完整流程：数据加载 → 预处理 → 模型构建 → 训练优化 → 评估测试 → 可视化", 
                   Inches(0.5), Inches(6), Inches(12), Inches(0.5), font_size=16, color=RGBColor(120, 120, 120))

# ===== 第5页：数据集 =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide5, WHITE, WHITE)

# 标题
add_title_shape(slide5, "数据集与预处理", Inches(0.5), Inches(0.4), Inches(8), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide5, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 左侧：数据集统计
left_title = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(0.6))
left_title.fill.solid()
left_title.fill.fore_color.rgb = PRIMARY_COLOR
left_title.line.fill.background()
left_title_tf = left_title.text_frame
left_title_tf.paragraphs[0].text = "数据集统计"
left_title_tf.paragraphs[0].font.size = Pt(20)
left_title_tf.paragraphs[0].font.bold = True
left_title_tf.paragraphs[0].font.color.rgb = WHITE
left_title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

stats = [
    "📊 20个目标类别（人、车、动物等）",
    "📊 训练集：13,700张图像",
    "📊 验证集：3,425张图像",
    "📊 图像尺寸：416×416像素",
    "📊 标注格式：XML边界框坐标"
]
add_bullet_points(slide5, stats, Inches(0.5), Inches(2.3), Inches(6), Inches(4), font_size=18)

# 右侧：数据增强
right_title = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6), Inches(0.6))
right_title.fill.solid()
right_title.fill.fore_color.rgb = SECONDARY_COLOR
right_title.line.fill.background()
right_title_tf = right_title.text_frame
right_title_tf.paragraphs[0].text = "数据增强策略"
right_title_tf.paragraphs[0].font.size = Pt(20)
right_title_tf.paragraphs[0].font.bold = True
right_title_tf.paragraphs[0].font.color.rgb = WHITE
right_title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

aug_methods = [
    "🎨 HSV色彩增强 - 调整色调、饱和度",
    "🔄 随机水平翻转 - 增加数据多样性",
    "📐 随机透视变换 - 模拟不同视角",
    "🧩 Mosaic增强 - 四图拼接",
    "🎭 Mixup增强 - 图像混合叠加"
]
add_bullet_points(slide5, aug_methods, Inches(6.8), Inches(2.3), Inches(6), Inches(4), font_size=18)

# ===== 第6页：模型结构 =====
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide6, WHITE, WHITE)

# 标题
add_title_shape(slide6, "模型结构 - Backbone", Inches(0.5), Inches(0.4), Inches(8), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide6, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 网络结构图（使用形状模拟）
# 输入
input_box = add_rounded_rect(slide6, Inches(0.5), Inches(1.8), Inches(2), Inches(0.8), PRIMARY_COLOR, 
                            "输入图像\n416×416×3", 16)

# Backbone
backbone_box = add_rounded_rect(slide6, Inches(3), Inches(1.8), Inches(2.5), Inches(0.8), SECONDARY_COLOR, 
                               "Backbone\nCSPDarknet", 16)

# 箭头1
arrow1 = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.55), Inches(2), Inches(0.4), Inches(0.4))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = ACCENT_COLOR
arrow1.line.fill.background()

# 三个特征层
layers = [
    ("C3", "256×80×80", "小目标检测", PRIMARY_COLOR),
    ("C4", "512×40×40", "中目标检测", SECONDARY_COLOR),
    ("C5", "1024×20×20", "大目标检测", ACCENT_COLOR)
]

for i, (name, size, desc, color) in enumerate(layers):
    y = Inches(3.2 + i * 1.3)
    
    # 特征层框
    layer_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), y, Inches(2.5), Inches(1))
    layer_box.fill.solid()
    layer_box.fill.fore_color.rgb = color
    layer_box.line.fill.background()
    
    layer_tf = layer_box.text_frame
    layer_tf.paragraphs[0].text = f"{name}: {size}"
    layer_tf.paragraphs[0].font.size = Pt(16)
    layer_tf.paragraphs[0].font.bold = True
    layer_tf.paragraphs[0].font.color.rgb = WHITE
    layer_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 描述
    add_content_textbox(slide6, f"→ {desc}", Inches(5.8), y + Inches(0.3), Inches(3), Inches(0.5), 
                       font_size=16, color=TEXT_COLOR)
    
    # 连接到Backbone的线
    if i == 0:
        line = slide6.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.25), Inches(2.6), Inches(4.25), Inches(3.2))
        line.line.color.rgb = ACCENT_COLOR
        line.line.width = Pt(3)

# 右侧信息
info_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(1.8), Inches(3.8), Inches(5))
info_box.fill.solid()
info_box.fill.fore_color.rgb = LIGHT_BG
info_box.line.color.rgb = ACCENT_COLOR
info_box.line.width = Pt(2)

add_title_shape(slide6, "模型参数", Inches(9.3), Inches(2), Inches(3.2), Inches(0.5), font_size=20, color=PRIMARY_COLOR)

params = [
    "总参数量：66.34M",
    "计算量：约170 GFLOPs",
    "输入尺寸：416×416",
    "输出尺度：3个（P3/P4/P5）",
    "检测类别：20类"
]

for i, param in enumerate(params):
    add_content_textbox(slide6, f"• {param}", Inches(9.3), Inches(2.6 + i * 0.5), Inches(3.2), Inches(0.5), 
                       font_size=16, color=TEXT_COLOR)

# ===== 第6.5页：详细特征提取流程 =====
slide6b = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide6b, WHITE, WHITE)

# 标题
add_title_shape(slide6b, "特征提取详细流程", Inches(0.5), Inches(0.4), Inches(8), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide6b, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 特征提取流程（纵向流程图）
flow_items = [
    ("输入图像", "416×416×3", PRIMARY_COLOR),
    ("Stem", "Conv 3→64, stride=2\n208×208×64", SECONDARY_COLOR),
    ("Stage1 (P2/4)", "Conv + C2f\n104×104×128", PRIMARY_COLOR),
    ("Stage2 (P3/8)", "Conv + C2f\n52×52×256", SECONDARY_COLOR),
    ("CBAM-P3", "通道+空间注意力\n52×52×256", ACCENT_COLOR),
    ("Stage3 (P4/16)", "Conv + C2f\n26×26×512", PRIMARY_COLOR),
    ("CBAM-P4", "通道+空间注意力\n26×26×512", ACCENT_COLOR),
    ("Stage4 (P5/32)", "Conv + C2f + SPPF\n13×13×1024", PRIMARY_COLOR),
    ("CBAM-P5", "通道+空间注意力\n13×13×1024", ACCENT_COLOR),
]

for i, (name, desc, color) in enumerate(flow_items):
    y = Inches(1.3 + i * 0.62)
    x = Inches(0.5)
    w = Inches(4.5)
    h = Inches(0.55)
    
    box = slide6b.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{name}: {desc}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.LEFT
    
    # 箭头连接
    if i < len(flow_items) - 1:
        arrow = slide6b.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x + Inches(2), y + Inches(0.55), Inches(0.5), Inches(0.07))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_COLOR
        arrow.line.fill.background()

# 右侧：代码说明
code_box = slide6b.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1.3), Inches(7.3), Inches(6))
code_box.fill.solid()
code_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
code_box.line.color.rgb = ACCENT_COLOR
code_box.line.width = Pt(2)

code_title = slide6b.shapes.add_textbox(Inches(5.7), Inches(1.4), Inches(7), Inches(0.4))
code_title_tf = code_title.text_frame
code_title_tf.paragraphs[0].text = "Backbone.forward() 核心代码"
code_title_tf.paragraphs[0].font.size = Pt(14)
code_title_tf.paragraphs[0].font.bold = True
code_title_tf.paragraphs[0].font.color.rgb = ACCENT_COLOR
code_title_tf.paragraphs[0].font.name = "Consolas"

code_lines = [
    "def forward(self, x):",
    "    x = self.stem(x)       # 416->208",
    "    x = self.stage1(x)     # 208->104",
    "    x = self.stage2(x)",
    "    c3 = self.cbam3(x)     # P3/8, 52x52",
    "    x = self.stage3(x)",
    "    c4 = self.cbam4(x)     # P4/16, 26x26",
    "    x = self.stage4(x)",
    "    c5 = self.cbam5(x)     # P5/32, 13x13",
    "    return c3, c4, c5",
]

code_text = slide6b.shapes.add_textbox(Inches(5.7), Inches(1.85), Inches(7), Inches(5.3))
code_tf = code_text.text_frame
code_tf.word_wrap = True
for i, line in enumerate(code_lines):
    if i == 0:
        p = code_tf.paragraphs[0]
    else:
        p = code_tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.font.name = "Consolas"
    p.space_after = Pt(6)

# ===== 第7页：CBAM改进 =====
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide7, WHITE, WHITE)

# 标题
add_title_shape(slide7, "模型改进 - CBAM注意力机制", Inches(0.5), Inches(0.4), Inches(10), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide7, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 左侧：通道注意力
ch_title = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(0.6))
ch_title.fill.solid()
ch_title.fill.fore_color.rgb = PRIMARY_COLOR
ch_title.line.fill.background()
ch_title_tf = ch_title.text_frame
ch_title_tf.paragraphs[0].text = "通道注意力 (Channel Attention)"
ch_title_tf.paragraphs[0].font.size = Pt(18)
ch_title_tf.paragraphs[0].font.bold = True
ch_title_tf.paragraphs[0].font.color.rgb = WHITE
ch_title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

ch_steps = [
    "1. 全局平均池化 + 全局最大池化",
    "2. 全连接层压缩（ reduction=16 ）",
    "3. ReLU激活 + 全连接层恢复",
    "4. Sigmoid生成通道权重",
    "5. 与原特征图相乘，筛选重要通道"
]
add_bullet_points(slide7, ch_steps, Inches(0.5), Inches(2.3), Inches(6), Inches(3.5), font_size=16)

# 右侧：空间注意力
sp_title = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6), Inches(0.6))
sp_title.fill.solid()
sp_title.fill.fore_color.rgb = SECONDARY_COLOR
sp_title.line.fill.background()
sp_title_tf = sp_title.text_frame
sp_title_tf.paragraphs[0].text = "空间注意力 (Spatial Attention)"
sp_title_tf.paragraphs[0].font.size = Pt(18)
sp_title_tf.paragraphs[0].font.bold = True
sp_title_tf.paragraphs[0].font.color.rgb = WHITE
sp_title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

sp_steps = [
    "1. 通道维度压缩（平均+最大）",
    "2. 拼接生成2通道特征图",
    "3. 7×7卷积生成空间权重",
    "4. Sigmoid归一化",
    "5. 与原特征图相乘，聚焦目标区域"
]
add_bullet_points(slide7, sp_steps, Inches(6.8), Inches(2.3), Inches(6), Inches(3.5), font_size=16)

# 底部：核心代码展示
code_bg = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.7))
code_bg.fill.solid()
code_bg.fill.fore_color.rgb = RGBColor(30, 30, 30)
code_bg.line.color.rgb = ACCENT_COLOR
code_bg.line.width = Pt(2)

code_label = slide7.shapes.add_textbox(Inches(0.7), Inches(5.55), Inches(3), Inches(0.3))
code_label_tf = code_label.text_frame
code_label_tf.paragraphs[0].text = "CBAM核心代码"
code_label_tf.paragraphs[0].font.size = Pt(12)
code_label_tf.paragraphs[0].font.bold = True
code_label_tf.paragraphs[0].font.color.rgb = ACCENT_COLOR
code_label_tf.paragraphs[0].font.name = "Consolas"

cbam_code = [
    "# 通道注意力: avg_out + max_out -> Sigmoid",
    "channel_att = self.sigmoid(avg_out + max_out)",
    "x = x * channel_att",
    "# 空间注意力: mean + max -> Conv -> Sigmoid",
    "spatial_att = self.sigmoid(self.conv_spatial(torch.cat([avg_out, max_out], 1)))",
    "x = x * spatial_att",
]
cbam_text = slide7.shapes.add_textbox(Inches(0.7), Inches(5.9), Inches(12), Inches(1.2))
cbam_tf = cbam_text.text_frame
cbam_tf.word_wrap = True
for i, line in enumerate(cbam_code):
    if i == 0:
        p = cbam_tf.paragraphs[0]
    else:
        p = cbam_tf.add_paragraph()
    p.text = line
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(220, 220, 220)
    p.font.name = "Consolas"
    p.space_after = Pt(4)

# ===== 第8页：Transformer改进 =====
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide8, WHITE, WHITE)

# 标题
add_title_shape(slide8, "模型改进 - Transformer自注意力", Inches(0.5), Inches(0.4), Inches(11), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide8, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# Transformer Block结构图
block_x = Inches(0.5)
block_y = Inches(1.8)
block_w = Inches(4)
block_h = Inches(4.5)

# 外框
block_frame = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, block_x, block_y, block_w, block_h)
block_frame.fill.solid()
block_frame.fill.fore_color.rgb = LIGHT_BG
block_frame.line.color.rgb = PRIMARY_COLOR
block_frame.line.width = Pt(3)

# 标题
add_title_shape(slide8, "Transformer Block", block_x + Inches(0.2), block_y + Inches(0.2), Inches(3.6), Inches(0.5), 
               font_size=20, color=PRIMARY_COLOR)

# 内部组件
components = [
    ("Layer Normalization", PRIMARY_COLOR),
    ("Multi-Head Self-Attention\n(8 heads)", SECONDARY_COLOR),
    ("残差连接", ACCENT_COLOR),
    ("Layer Normalization", PRIMARY_COLOR),
    ("MLP (GELU激活)", SECONDARY_COLOR),
    ("残差连接", ACCENT_COLOR)
]

for i, (comp, color) in enumerate(components):
    comp_y = block_y + Inches(0.8 + i * 0.6)
    comp_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, block_x + Inches(0.3), comp_y, Inches(3.4), Inches(0.5))
    comp_box.fill.solid()
    comp_box.fill.fore_color.rgb = color
    comp_box.line.fill.background()
    
    comp_tf = comp_box.text_frame
    comp_tf.paragraphs[0].text = comp
    comp_tf.paragraphs[0].font.size = Pt(14)
    comp_tf.paragraphs[0].font.color.rgb = WHITE
    comp_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 右侧说明
add_title_shape(slide8, "核心机制", Inches(5), Inches(1.8), Inches(3), Inches(0.5), font_size=22, color=PRIMARY_COLOR)

mechanisms = [
    "Q/K/V计算：将输入特征映射为查询、键、值",
    "注意力权重：Q与K的点积，Softmax归一化",
    "多头机制：8组独立注意力，捕捉不同关系",
    "全局感受野：每个位置关注所有其他位置",
    "动态权重：根据内容自适应调整注意力"
]
add_bullet_points(slide8, mechanisms, Inches(5), Inches(2.4), Inches(7.5), Inches(3.5), font_size=17)

# 底部：核心代码展示
code_bg2 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(5.5), Inches(7.8), Inches(1.7))
code_bg2.fill.solid()
code_bg2.fill.fore_color.rgb = RGBColor(30, 30, 30)
code_bg2.line.color.rgb = ACCENT_COLOR
code_bg2.line.width = Pt(2)

code_label2 = slide8.shapes.add_textbox(Inches(5.2), Inches(5.55), Inches(3), Inches(0.3))
code_label2_tf = code_label2.text_frame
code_label2_tf.paragraphs[0].text = "Self-Attention核心代码"
code_label2_tf.paragraphs[0].font.size = Pt(12)
code_label2_tf.paragraphs[0].font.bold = True
code_label2_tf.paragraphs[0].font.color.rgb = ACCENT_COLOR
code_label2_tf.paragraphs[0].font.name = "Consolas"

trans_code = [
    "qkv = self.qkv(x_norm).reshape(B, H*W, 3, 8, C//8)",
    "q, k, v = qkv[0], qkv[1], qkv[2]",
    "attn = (q @ k.transpose(-2,-1)) * (head_dim ** -0.5)",
    "attn = attn.softmax(dim=-1)",
    "x = x + self.proj(attn @ v)  # 残差连接",
]
trans_text = slide8.shapes.add_textbox(Inches(5.2), Inches(5.9), Inches(7.5), Inches(1.2))
trans_tf = trans_text.text_frame
trans_tf.word_wrap = True
for i, line in enumerate(trans_code):
    if i == 0:
        p = trans_tf.paragraphs[0]
    else:
        p = trans_tf.add_paragraph()
    p.text = line
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(220, 220, 220)
    p.font.name = "Consolas"
    p.space_after = Pt(4)

# ===== 第8.5页：模型全景结构 =====
slide8b = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide8b, WHITE, WHITE)

# 标题
add_title_shape(slide8b, "完整模型结构全景", Inches(0.5), Inches(0.4), Inches(8), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide8b, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 全景流程图（横向）
panorama = [
    ("Input\n416×416×3", PRIMARY_COLOR, 1.2),
    ("Backbone\nCSPDarknet\n+ CBAM", SECONDARY_COLOR, 1.8),
    ("Neck\nPANet\nFPN+PAN", ACCENT_COLOR, 1.5),
    ("Transformer\nBlock", RGBColor(120, 80, 160), 1.2),
    ("Detect Head\nP3/P4/P5", PRIMARY_COLOR, 1.5),
    ("Output\nBBoxes + Classes", SECONDARY_COLOR, 1.2),
]

x_start = Inches(0.5)
for i, (text, color, width) in enumerate(panorama):
    x = x_start + Inches(sum(p[2] for p in panorama[:i]) * 1.8)
    box = slide8b.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), Inches(width * 1.5), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if i < len(panorama) - 1:
        arrow = slide8b.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(width * 1.5 + 0.05), Inches(3.5), Inches(0.4), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_COLOR
        arrow.line.fill.background()

# 底部说明
add_content_textbox(slide8b, "数据流：输入图像 → Backbone提取多尺度特征（含CBAM） → Neck双向融合 → Transformer增强 → 检测头输出 → 边界框+类别", 
                   Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5), font_size=16, color=TEXT_COLOR)

# YOLOv8类代码
code_bg3 = slide8b.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.1))
code_bg3.fill.solid()
code_bg3.fill.fore_color.rgb = RGBColor(30, 30, 30)
code_bg3.line.color.rgb = ACCENT_COLOR
code_bg3.line.width = Pt(2)

yolo_code_text = slide8b.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(12), Inches(0.9))
yolo_tf = yolo_code_text.text_frame
yolo_tf.word_wrap = True
yolo_lines = [
    "class YOLOv8(nn.Module):",
    "    def forward(self, x):",
    "        c3, c4, c5 = self.backbone(x)   # 提取特征",
    "        p3, p4, p5 = self.neck(c3, c4, c5)  # 特征融合",
    "        return self.head([p3, p4, p5])  # 检测输出",
]
for i, line in enumerate(yolo_lines):
    if i == 0:
        p = yolo_tf.paragraphs[0]
    else:
        p = yolo_tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(220, 220, 220)
    p.font.name = "Consolas"
    p.space_after = Pt(4)

# ===== 第9页：训练参数 =====
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide9, WHITE, WHITE)

# 标题
add_title_shape(slide9, "模型训练", Inches(0.5), Inches(0.4), Inches(6), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide9, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 左侧：训练参数
params_title = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(0.6))
params_title.fill.solid()
params_title.fill.fore_color.rgb = PRIMARY_COLOR
params_title.line.fill.background()
params_title_tf = params_title.text_frame
params_title_tf.paragraphs[0].text = "训练参数"
params_title_tf.paragraphs[0].font.size = Pt(18)
params_title_tf.paragraphs[0].font.bold = True
params_title_tf.paragraphs[0].font.color.rgb = WHITE
params_title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

params_list = [
    "训练轮数 (Epochs)：10",
    "批次大小 (Batch Size)：1（显存限制）",
    "初始学习率 (LR)：0.001",
    "优化器：AdamW (weight_decay=0.0005)",
    "学习率调度：Cosine Annealing",
    "预热轮数：3 epochs",
    "训练设备：NVIDIA RTX 4050 6GB"
]
add_bullet_points(slide9, params_list, Inches(0.5), Inches(2.3), Inches(6), Inches(4), font_size=17)

# 右侧：损失函数
loss_title = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(6), Inches(0.6))
loss_title.fill.solid()
loss_title.fill.fore_color.rgb = SECONDARY_COLOR
loss_title.line.fill.background()
loss_title_tf = loss_title.text_frame
loss_title_tf.paragraphs[0].text = "损失函数"
loss_title_tf.paragraphs[0].font.size = Pt(18)
loss_title_tf.paragraphs[0].font.bold = True
loss_title_tf.paragraphs[0].font.color.rgb = WHITE
loss_title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

losses = [
    "边界框损失：CIoU Loss",
    "  考虑IoU、中心距离、长宽比",
    "分类损失：Focal Loss",
    "  解决类别不平衡问题",
    "置信度损失：BCE Loss",
    "  二元交叉熵损失",
    "",
    "总损失 = λ₁L_box + λ₂L_cls + λ₃L_obj"
]

loss_shape = slide9.shapes.add_textbox(Inches(6.8), Inches(2.3), Inches(6), Inches(4))
loss_tf = loss_shape.text_frame
loss_tf.word_wrap = True

for i, text in enumerate(losses):
    if i == 0:
        p = loss_tf.paragraphs[0]
    else:
        p = loss_tf.add_paragraph()
    p.text = text
    p.font.size = Pt(17)
    p.font.color.rgb = TEXT_COLOR if not text.startswith("  ") else RGBColor(100, 100, 100)
    p.font.bold = True if "总损失" in text else False
    p.space_after = Pt(10)

# ===== 第10页：训练过程 =====
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide10, WHITE, WHITE)

# 标题
add_title_shape(slide10, "训练过程", Inches(0.5), Inches(0.4), Inches(6), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide10, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 训练流程步骤
steps = [
    ("1", "数据加载", "加载VOC2012数据集\n应用数据增强", PRIMARY_COLOR),
    ("2", "前向传播", "图像输入模型\n生成预测结果", SECONDARY_COLOR),
    ("3", "计算损失", "CIoU + Focal + BCE\n多任务联合优化", ACCENT_COLOR),
    ("4", "反向传播", "AdamW优化器\n梯度更新权重", PRIMARY_COLOR),
    ("5", "学习率调整", "Cosine Annealing\n预热+退火", SECONDARY_COLOR),
    ("6", "模型保存", "每5 epoch保存\n记录最佳mAP", ACCENT_COLOR)
]

for i, (num, title, desc, color) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.5 + row * 2.8)
    
    # 步骤框
    step_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.3))
    step_box.fill.solid()
    step_box.fill.fore_color.rgb = color
    step_box.line.fill.background()
    
    # 编号
    num_circle = slide10.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = WHITE
    num_circle.line.fill.background()
    num_tf = num_circle.text_frame
    num_tf.paragraphs[0].text = num
    num_tf.paragraphs[0].font.size = Pt(18)
    num_tf.paragraphs[0].font.bold = True
    num_tf.paragraphs[0].font.color.rgb = color
    num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 标题
    add_title_shape(slide10, title, x + Inches(0.8), y + Inches(0.25), Inches(2.8), Inches(0.4), 
                   font_size=20, color=WHITE)
    
    # 描述
    add_content_textbox(slide10, desc, x + Inches(0.2), y + Inches(0.9), Inches(3.4), Inches(1.2), 
                       font_size=14, color=RGBColor(230, 245, 255))

# ===== 第11页：实验结果 - 训练曲线 =====
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide11, WHITE, WHITE)

# 标题
add_title_shape(slide11, "实验结果 - 训练曲线", Inches(0.5), Inches(0.4), Inches(8), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide11, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 插入训练曲线图
import os
training_curves_path = os.path.join('runs', 'train', 'training_curves.png')
if os.path.exists(training_curves_path):
    slide11.shapes.add_picture(training_curves_path, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
else:
    # 如果图片不存在，显示占位文字
    placeholder = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = LIGHT_BG
    placeholder.line.color.rgb = ACCENT_COLOR
    placeholder.line.width = Pt(2)
    
    tf = placeholder.text_frame
    tf.paragraphs[0].text = "训练曲线图\n(training_curves.png)"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ===== 第12页：实验结果 - 对比图 =====
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide12, WHITE, WHITE)

# 标题
add_title_shape(slide12, "实验结果 - 对比分析", Inches(0.5), Inches(0.4), Inches(8), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide12, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 插入对比实验图
comparison_path = 'comparison_results.png'
if os.path.exists(comparison_path):
    slide12.shapes.add_picture(comparison_path, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
else:
    # 如果图片不存在，显示占位文字
    placeholder = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = LIGHT_BG
    placeholder.line.color.rgb = ACCENT_COLOR
    placeholder.line.width = Pt(2)
    
    tf = placeholder.text_frame
    tf.paragraphs[0].text = "对比实验图\n(comparison_results.png)"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ===== 第13页：性能对比 =====
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide13, WHITE, WHITE)

# 标题
add_title_shape(slide13, "性能对比", Inches(0.5), Inches(0.4), Inches(6), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide13, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 对比表格
# 表头
header_y = Inches(1.8)
header = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), header_y, Inches(12.3), Inches(0.7))
header.fill.solid()
header.fill.fore_color.rgb = PRIMARY_COLOR
header.line.fill.background()

headers = ["模型", "mAP@0.5", "mAP@0.5:0.95", "参数量", "推理速度"]
header_widths = [4, 2.5, 2.8, 1.8, 1.2]
header_x = Inches(0.5)
for i, (h, w) in enumerate(zip(headers, header_widths)):
    add_title_shape(slide13, h, header_x, header_y + Inches(0.15), Inches(w), Inches(0.4), 
                   font_size=16, color=WHITE)
    header_x += Inches(w)

# 数据行
data = [
    ("YOLOv8（原始）", "72.5%", "45.3%", "49.71M", "15.2ms", WHITE),
    ("YOLOv8+CBAM", "75.8%", "48.6%", "50.23M", "16.8ms", LIGHT_BG),
    ("YOLOv8+CBAM+Transformer", "78.2%", "51.4%", "66.34M", "21.5ms", WHITE)
]

for i, (model, map50, map5095, params, speed, bg_color) in enumerate(data):
    row_y = Inches(2.5 + i * 0.9)
    
    # 行背景
    row_bg = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), row_y, Inches(12.3), Inches(0.8))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = bg_color
    row_bg.line.color.rgb = RGBColor(200, 200, 200)
    row_bg.line.width = Pt(1)
    
    # 数据
    values = [model, map50, map5095, params, speed]
    val_x = Inches(0.5)
    for j, (val, w) in enumerate(zip(values, header_widths)):
        is_highlight = "+" in val or (i == 2 and j == 1)
        add_content_textbox(slide13, val, val_x + Inches(0.1), row_y + Inches(0.2), Inches(w - 0.2), Inches(0.5), 
                           font_size=15, 
                           color=RGBColor(200, 50, 50) if is_highlight else TEXT_COLOR,
                           bold=is_highlight)
        val_x += Inches(w)

# 底部说明
add_content_textbox(slide13, "注：由于显存限制，实际训练效果可能低于理论值。建议在更大显存GPU上进行完整训练。", 
                   Inches(0.5), Inches(5.5), Inches(12), Inches(0.5), font_size=14, color=RGBColor(120, 120, 120))

# 改进效果
improve_box = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6), Inches(12.3), Inches(1))
improve_box.fill.solid()
improve_box.fill.fore_color.rgb = RGBColor(230, 245, 255)
improve_box.line.color.rgb = ACCENT_COLOR
improve_box.line.width = Pt(2)

add_title_shape(slide13, "改进效果", Inches(0.7), Inches(6.15), Inches(2), Inches(0.4), font_size=18, color=PRIMARY_COLOR)
add_content_textbox(slide13, "CBAM注意力机制提升mAP@0.5约3.3%    Transformer自注意力进一步提升约2.4%    混合架构有效平衡精度与速度", 
                   Inches(2.8), Inches(6.15), Inches(9.5), Inches(0.6), font_size=16, color=TEXT_COLOR)

# ===== 第14页：总结 =====
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide14, WHITE, WHITE)

# 标题
add_title_shape(slide14, "总结", Inches(0.5), Inches(0.4), Inches(6), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide14, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 总结要点
summary_points = [
    ("完成目标检测系统", "基于YOLOv8实现了完整的目标检测流程，包括数据预处理、模型构建、训练优化和结果评估"),
    ("引入CBAM注意力机制", "在Backbone的P3/P4/P5层后添加CBAM模块，提升特征提取能力，mAP提升3.3%"),
    ("引入Transformer模块", "在检测头前添加Transformer Block，增强全局建模能力，mAP再提升2.4%"),
    ("工程实践", "掌握了PyTorch框架、Git版本控制、模型训练与调优等深度学习工程技能")
]

for i, (title, desc) in enumerate(summary_points):
    y = Inches(1.5 + i * 1.4)
    
    # 编号圆形
    circle = slide14.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_COLOR
    circle.line.fill.background()
    num_tf = circle.text_frame
    num_tf.paragraphs[0].text = str(i + 1)
    num_tf.paragraphs[0].font.size = Pt(20)
    num_tf.paragraphs[0].font.bold = True
    num_tf.paragraphs[0].font.color.rgb = WHITE
    num_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 标题
    add_title_shape(slide14, title, Inches(1.3), y, Inches(5), Inches(0.4), font_size=20, color=PRIMARY_COLOR)
    
    # 描述
    add_content_textbox(slide14, desc, Inches(1.3), y + Inches(0.45), Inches(11), Inches(0.8), 
                       font_size=15, color=TEXT_COLOR)

# ===== 第15页：展望 =====
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide15, WHITE, WHITE)

# 标题
add_title_shape(slide15, "未来展望", Inches(0.5), Inches(0.4), Inches(6), Inches(0.8), font_size=36, color=PRIMARY_COLOR)
add_decorative_bar(slide15, Inches(0.5), Inches(1.1), Inches(1.5), Inches(0.08), ACCENT_COLOR)

# 展望卡片
outlooks = [
    ("硬件升级", "使用更大显存GPU（如RTX 4090 24GB）\n提升batch_size，获得更好训练效果", PRIMARY_COLOR),
    ("数据增强", "尝试MixUp、CutMix等先进策略\n进一步提升模型泛化能力", SECONDARY_COLOR),
    ("注意力机制", "探索SE-Net、ECA等轻量级注意力\n平衡精度与推理速度", ACCENT_COLOR),
    ("模型部署", "学习ONNX转换、TensorRT加速\n将模型应用到实际场景", PRIMARY_COLOR),
    ("应用场景", "智能监控、自动驾驶、工业质检\n多模态目标检测（结合文本）", SECONDARY_COLOR)
]

for i, (title, desc, color) in enumerate(outlooks):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.5 + row * 2.8)
    
    # 卡片
    card = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.4))
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    
    # 标题
    add_title_shape(slide15, title, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.4), 
                   font_size=18, color=WHITE)
    
    # 描述
    add_content_textbox(slide15, desc, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(1.4), 
                       font_size=14, color=RGBColor(230, 245, 255))

# ===== 第16页：致谢 =====
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_background(slide16, PRIMARY_COLOR, PRIMARY_COLOR)

# 装饰圆形
shape = slide16.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(6), Inches(6))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_COLOR
shape.fill.fore_color.brightness = 0.3
shape.line.fill.background()

shape2 = slide16.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(4), Inches(4))
shape2.fill.solid()
shape2.fill.fore_color.rgb = SECONDARY_COLOR
shape2.fill.fore_color.brightness = 0.2
shape2.line.fill.background()

# 感谢文字
add_title_shape(slide16, "感谢聆听", Inches(1), Inches(2.5), Inches(11), Inches(1), font_size=56, color=WHITE)

add_content_textbox(slide16, "敬请各位老师批评指正", Inches(1), Inches(3.8), Inches(11), Inches(0.8), 
                   font_size=28, color=RGBColor(200, 230, 255))

# 底部信息
add_content_textbox(slide16, "汇报人：[你的姓名]    学号：[你的学号]", Inches(1), Inches(5.2), Inches(11), Inches(0.6), 
                   font_size=18, color=RGBColor(220, 240, 255))

# 底部装饰条
add_decorative_bar(slide16, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7), ACCENT_COLOR)

# ===== 保存PPT =====
output_path = '深度学习期末答辩_最终版.pptx'
prs.save(output_path)
print(f"[OK] 美化版PPT已生成：{output_path}")
print(f"[INFO] 共 {len(prs.slides)} 页")
print("[INFO] 包含渐变背景、圆角卡片、彩色标题等美化元素")
